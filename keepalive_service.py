#!/usr/bin/env python3
"""
Keepalive Service — 保活层
职责: iLink长轮询 + 暗号秒回 + 欢迎推送 + session保活
不调Hermes，不处理AI逻辑，不依赖网关。

非暗号消息 → POST 到 Agent Connector (:9101) → 等回复 → 推回iLink
"""

import asyncio
import base64
import hashlib
import json
import logging
import mimetypes
import os
import re
import secrets
import signal
import struct
import sys
import time
import uuid
from dataclasses import dataclass
from datetime import date
from pathlib import Path
from typing import Optional
from urllib.parse import quote

import aiohttp

# ══════════════════════════════════════════════
# Config
# ══════════════════════════════════════════════

# iLink
ILINK_BASE_URL = "https://ilinkai.weixin.qq.com"
WEIXIN_CDN_BASE_URL = "https://novac2c.cdn.weixin.qq.com/c2c"
ILINK_APP_ID = "bot"
ILINK_APP_CLIENT_VERSION = (2 << 16) | (2 << 8) | 0  # 131584

LONG_POLL_TIMEOUT_MS = 35_000
SHORT_POLL_TIMEOUT_MS = 5_000
API_TIMEOUT_MS = 15_000
SESSION_EXPIRED_ERRCODE = -14
SESSION_PAUSE_MS = 10 * 60 * 1000  # 10min
HEARTBEAT_INTERVAL_MS = 5 * 60 * 1000  # 5min

# Agent Connector
AGENT_URL = "http://127.0.0.1:9101/api/message"
AGENT_TIMEOUT = 120  # Hermes 最长回复时间


# HTTP API (本服务)
HTTP_PORT = 9100

# DB
DB_DSN = "postgresql://lucky:lucky_pass@localhost:5432/weclawd"
DB_RELOAD_INTERVAL = 30  # 30s

# Paths
STATE_DIR = Path.home() / ".hermes" / "keepalive"
SYNC_BUF_DIR = STATE_DIR / "sync_buf"
PAUSE_DIR = STATE_DIR / "pause"
TOKEN_STORE_DIR = STATE_DIR / "tokens"
CONFIG_DIR = Path.home() / "weclaw-1" / "config"
CODE_PHRASE_FILE = CONFIG_DIR / "access_codes.json"

STATE_DIR.mkdir(parents=True, exist_ok=True)
SYNC_BUF_DIR.mkdir(parents=True, exist_ok=True)
PAUSE_DIR.mkdir(parents=True, exist_ok=True)
TOKEN_STORE_DIR.mkdir(parents=True, exist_ok=True)
CONFIG_DIR.mkdir(parents=True, exist_ok=True)

# Message types
MSG_TYPE_BOT = 2
MSG_STATE_FINISH = 2
ITEM_TEXT = 1
ITEM_IMAGE = 2
ITEM_VOICE = 3
ITEM_FILE = 4
ITEM_VIDEO = 5

TYPING_START = 1
TYPING_STOP = 2

# ══════════════════════════════════════════════
# 暗号管理
# ══════════════════════════════════════════════

_DEFAULT_CODES = {
    "codes": [
{"code": "天王盖地虎", "reply": "OpenClaw 是SB！", "level": "admin"},
        {"code": "宝塔镇河妖", "reply": "微侠真牛逼！", "level": "admin"},
        {"code": "微侠真牛逼", "reply": "天王盖地虎！同志！", "level": "admin"},
        {"code": "openclaw是sb", "reply": "宝塔镇河妖！收到！", "level": "admin"},
    ]
}

def _ensure_code_file():
    if not CODE_PHRASE_FILE.exists():
        CODE_PHRASE_FILE.write_text(json.dumps(_DEFAULT_CODES, ensure_ascii=False, indent=2))

def load_code_phrases() -> list[dict]:
    _ensure_code_file()
    try:
        data = json.loads(CODE_PHRASE_FILE.read_text())
        return data.get("codes", [])
    except (json.JSONDecodeError, FileNotFoundError):
        return _DEFAULT_CODES["codes"]

def match_code_phrase(text: str) -> Optional[dict]:
    stripped = re.sub(r'[，。！？、；：""\'\s]', '', text).lower()
    codes = load_code_phrases()
    for entry in codes:
        code = re.sub(r'[，。！？、；：""\'\s]', '', entry["code"]).lower()
        if code in stripped:
            return entry
    return None

# ══════════════════════════════════════════════
# 工具函数
# ══════════════════════════════════════════════

def _build_base_info() -> dict:
    return {"channel_version": "2.2.0"}

def _build_headers(token: str, body: str = "") -> dict:
    headers = {
        "Content-Type": "application/json",
        "AuthorizationType": "ilink_bot_token",
        "iLink-App-Id": ILINK_APP_ID,
        "iLink-App-ClientVersion": str(ILINK_APP_CLIENT_VERSION),
    }
    if body:
        headers["Content-Length"] = str(len(body.encode("utf-8")))
        uin_val = struct.unpack(">I", secrets.token_bytes(4))[0]
        headers["X-WECHAT-UIN"] = base64.b64encode(str(uin_val).encode("utf-8")).decode("ascii")
    if token:
        headers["Authorization"] = f"Bearer {token}"
    return headers

def _pkcs7_pad(data: bytes, block_size: int = 16) -> bytes:
    pad_len = block_size - (len(data) % block_size)
    return data + bytes([pad_len] * pad_len)

def _aes128_ecb_encrypt(plaintext: bytes, key: bytes) -> bytes:
    from cryptography.hazmat.primitives.ciphers import Cipher, algorithms, modes
    cipher = Cipher(algorithms.AES128(key), modes.ECB())
    encryptor = cipher.encryptor()
    return encryptor.update(_pkcs7_pad(plaintext)) + encryptor.finalize()

def _aes_padded_size(size: int) -> int:
    block = 16
    return ((size + block - 1) // block) * block

def _safe_key(s: str) -> str:
    return s.replace("@", "_at_").replace(".", "_dot_")

# ══════════════════════════════════════════════
# TokenStore — disk-backed context_token persistence
# ══════════════════════════════════════════════

class TokenStore:
    """Persist context_tokens per (bot_id, user_id) so proactive
    pushes (cron, CS forwarding) can send without an inbound message."""

    def _tokens_file(self, bot_id: str) -> Path:
        return TOKEN_STORE_DIR / f"{_safe_key(bot_id)}.json"

    def get(self, bot_id: str, user_id: str) -> str:
        """Get saved context_token for a user. Returns '' if none."""
        try:
            data = json.loads(self._tokens_file(bot_id).read_text())
            return data.get(user_id, "")
        except (FileNotFoundError, json.JSONDecodeError):
            return ""

    def set(self, bot_id: str, user_id: str, token: str):
        """Save context_token for a user."""
        f = self._tokens_file(bot_id)
        try:
            data = json.loads(f.read_text())
        except (FileNotFoundError, json.JSONDecodeError):
            data = {}
        if token:
            data[user_id] = token
        else:
            data.pop(user_id, None)
        f.write_text(json.dumps(data, ensure_ascii=False))

    def remove(self, bot_id: str, user_id: str):
        """Remove saved token (e.g. on session expiry)."""
        self.set(bot_id, user_id, "")

    def clear_bot(self, bot_id: str):
        """Remove all tokens for a bot."""
        f = self._tokens_file(bot_id)
        f.unlink(missing_ok=True)

_token_store = TokenStore()

def _sync_path(bot_id: str) -> Path:
    return SYNC_BUF_DIR / f"{_safe_key(bot_id)}.json"

def _pause_path(bot_id: str) -> Path:
    return PAUSE_DIR / f"{_safe_key(bot_id)}.json"

def _make_logger(name: str) -> logging.Logger:
    log = logging.getLogger(name)
    log.setLevel(logging.DEBUG)
    log.handlers.clear()
    fmt = logging.Formatter("%(asctime)s [%(name)s] %(levelname)s %(message)s", datefmt="%H:%M:%S")
    sh = logging.StreamHandler(sys.stderr)
    sh.setFormatter(fmt)
    log.addHandler(sh)
    return log

root_log = _make_logger("keepalive")

# Session persistence
def load_sync_buf(bot_id: str) -> str:
    try:
        return json.loads(_sync_path(bot_id).read_text()).get("buf", "")
    except (FileNotFoundError, json.JSONDecodeError):
        return ""

def save_sync_buf(bot_id: str, buf: str):
    _sync_path(bot_id).write_text(json.dumps({"buf": buf}))

def is_paused(bot_id: str) -> bool:
    try:
        until = json.loads(_pause_path(bot_id).read_text()).get("until", 0)
        if time.time() * 1000 < until:
            return True
        _pause_path(bot_id).unlink(missing_ok=True)
        return False
    except (FileNotFoundError, json.JSONDecodeError):
        return False

def pause_bot(bot_id: str, duration_ms: int = SESSION_PAUSE_MS):
    until = int(time.time() * 1000 + duration_ms)
    _pause_path(bot_id).write_text(json.dumps({"until": until}))

def clear_pause(bot_id: str):
    _pause_path(bot_id).unlink(missing_ok=True)

def pause_remaining_ms(bot_id: str) -> int:
    try:
        until = json.loads(_pause_path(bot_id).read_text()).get("until", 0)
        return max(0, until - int(time.time() * 1000))
    except (FileNotFoundError, json.JSONDecodeError):
        return 0

# ══════════════════════════════════════════════
# iLink API
# ══════════════════════════════════════════════

async def _ilink_post(endpoint: str, payload: dict, token: str,
                      timeout_ms: int = API_TIMEOUT_MS,
                      session: Optional[aiohttp.ClientSession] = None) -> dict:
    body = json.dumps({**payload, "base_info": _build_base_info()}, separators=(",", ":"))
    headers = _build_headers(token, body)

    async def _do(s: aiohttp.ClientSession) -> dict:
        async with s.post(f"{ILINK_BASE_URL}/{endpoint}",
                          data=body, headers=headers,
                          timeout=aiohttp.ClientTimeout(total=timeout_ms / 1000)) as r:
            if r.status != 200:
                raise RuntimeError(f"HTTP {r.status}: {await r.text()}")
            raw = await r.read()
            return json.loads(raw)

    if session:
        return await _do(session)
    async with aiohttp.ClientSession() as s:
        return await _do(s)

async def notify_start(token: str) -> bool:
    try:
        await _ilink_post("ilink/bot/msg/notifystart", {}, token, timeout_ms=10_000)
        return True
    except Exception:
        return False

async def notify_stop(token: str) -> bool:
    try:
        await _ilink_post("ilink/bot/msg/notifystop", {}, token, timeout_ms=5_000)
        return True
    except Exception:
        return False

async def get_updates(token: str, sync_buf: str, timeout_ms: int,
                      session: aiohttp.ClientSession) -> dict:
    body = json.dumps({
        "get_updates_buf": sync_buf,
        "base_info": _build_base_info(),
    }, separators=(",", ":"))
    headers = _build_headers(token, body)
    try:
        async with session.post(f"{ILINK_BASE_URL}/ilink/bot/getupdates",
                                data=body, headers=headers,
                                timeout=aiohttp.ClientTimeout(total=timeout_ms / 1000)) as r:
            if r.status != 200:
                return {"ret": -1, "errcode": r.status}
            raw = await r.read()
            return json.loads(raw)
    except asyncio.TimeoutError:
        return {"ret": 0, "msgs": [], "get_updates_buf": sync_buf}

async def _ilink_get_config(token: str, user_id: str,
                            context_token: Optional[str] = None,
                            session: Optional[aiohttp.ClientSession] = None) -> dict:
    payload = {"ilink_user_id": user_id}
    if context_token:
        payload["context_token"] = context_token
    return await _ilink_post("ilink/bot/getconfig", payload, token,
                             timeout_ms=10_000, session=session)

# ══════════════════════════════════════════════
# 消息发送
# ══════════════════════════════════════════════

def _aes_pkcs7_pad(data: bytes, block_size: int = 16) -> bytes:
    pad_len = block_size - (len(data) % block_size)
    return data + bytes([pad_len] * pad_len)

def _aes_padded_size(size: int) -> int:
    return ((size + 1 + 15) // 16) * 16

def _search_file(filename: str) -> str:
    """当 MEDIA 路径不存在时，在已知目录下搜索文件。返回第一个匹配的绝对路径。"""
    basename = os.path.basename(filename)
    search_dirs = [
        "/home/ubuntu/weclaw-keepalive/downloads",
        "/home/ubuntu/weclaw-1/app/static/downloads",
        "/mnt/shared/projects",
    ]
    for d in search_dirs:
        if os.path.isfile(d) and basename in d:
            return d
        if os.path.isdir(d):
            for root, dirs, files in os.walk(d):
                for f in files:
                    if f == basename or basename in f:
                        return os.path.join(root, f)
    return ""


async def _send_file(token: str, to_user_id: str, file_path: str,
                     caption: str = "", context_token: str = "",
                     bot_id: str = ""):
    """通过 iLink CDN 上传文件并推送媒体消息"""
    path = Path(file_path)
    if not path.exists():
        root_log.warning("[send_file] 文件不存在: %s", file_path)
        return False
    try:
        plaintext = path.read_bytes()
    except Exception as e:
        root_log.warning("[send_file] 读文件失败: %s", e)
        return False

    # 确定媒体类型
    ext = path.suffix.lower()
    mime = mimetypes.guess_type(str(path))[0] or "application/octet-stream"
    if mime.startswith("image/"):
        media_type = 1  # MEDIA_IMAGE
        item_type = ITEM_IMAGE
        item_key = "image_item"
        item_builder = lambda **kw: {"type": item_type, item_key: {
            "media": {"encrypt_query_param": kw["ep"], "aes_key": kw["ak"], "encrypt_type": 1},
            "mid_size": kw["cs"],
        }}
    elif mime.startswith("video/"):
        media_type = 2  # MEDIA_VIDEO
        item_type = ITEM_VIDEO
        item_key = "video_item"
        item_builder = lambda **kw: {"type": item_type, item_key: {
            "media": {"encrypt_query_param": kw["ep"], "aes_key": kw["ak"], "encrypt_type": 1},
            "video_size": kw["cs"], "play_length": 0, "video_md5": kw["md5"],
        }}
    elif ext in (".silk",):
        media_type = 4  # MEDIA_VOICE
        item_type = ITEM_VOICE
        item_key = "voice_item"
        item_builder = lambda **kw: {"type": item_type, item_key: {
            "media": {"encrypt_query_param": kw["ep"], "aes_key": kw["ak"], "encrypt_type": 1},
            "encode_type": 6, "sample_rate": 24000, "bits_per_sample": 16,
            "playtime": 0,
        }}
    else:
        media_type = 3  # MEDIA_FILE
        item_type = ITEM_FILE
        item_key = "file_item"
        item_builder = lambda **kw: {"type": item_type, item_key: {
            "media": {"encrypt_query_param": kw["ep"], "aes_key": kw["ak"], "encrypt_type": 1},
            "file_name": kw["fn"], "len": str(kw["ps"]),
        }}

    filekey = secrets.token_hex(16)
    aes_key = secrets.token_bytes(16)
    rawsize = len(plaintext)
    rawfilemd5 = hashlib.md5(plaintext).hexdigest()
    padded = _aes_pkcs7_pad(plaintext)
    filesize = len(padded)

    try:
        # 获取上传地址
        upload_resp = await _ilink_post("ilink/bot/getuploadurl", {
            "filekey": filekey,
            "media_type": media_type,
            "to_user_id": to_user_id,
            "rawsize": rawsize,
            "rawfilemd5": rawfilemd5,
            "filesize": filesize,
            "no_need_thumb": True,
            "aeskey": aes_key.hex(),
        }, token, timeout_ms=15_000)
    except Exception as e:
        root_log.warning("[send_file] getuploadurl 失败: %s", e)
        return False

    upload_param = upload_resp.get("upload_param") or upload_resp.get("encrypted_query_param", "")
    upload_full_url = upload_resp.get("upload_full_url", "")
    if not upload_param and not upload_full_url:
        root_log.warning("[send_file] 无上传地址: %s", upload_resp)
        return False

    upload_url = upload_full_url or (
        f"{WEIXIN_CDN_BASE_URL}/upload"
        f"?encrypted_query_param={quote(upload_param, safe='')}"
        f"&filekey={quote(filekey, safe='')}"
    )

    # AES-128-ECB 加密
    from cryptography.hazmat.primitives.ciphers import Cipher, algorithms, modes
    from cryptography.hazmat.backends import default_backend
    cipher = Cipher(algorithms.AES(aes_key), modes.ECB(), backend=default_backend())
    encryptor = cipher.encryptor()
    ciphertext = encryptor.update(padded) + encryptor.finalize()

    # 上传到 CDN
    try:
        async with aiohttp.ClientSession() as sess:
            async with sess.post(upload_url, data=ciphertext,
                                 headers={"Content-Type": "application/octet-stream"},
                                 timeout=aiohttp.ClientTimeout(total=30)) as r:
                if r.status != 200:
                    body = await r.read()
                    root_log.warning("[send_file] CDN 上传失败 HTTP %d: %s", r.status, body[:200])
                    return False
                encrypted_param = r.headers.get("x-encrypted-param", upload_param)
    except Exception as e:
        root_log.warning("[send_file] CDN 上传异常: %s", e)
        return False

    # 构造媒体消息
    aes_key_for_api = base64.b64encode(aes_key.hex().encode("ascii")).decode("ascii")
    kwargs = {"ep": encrypted_param, "ak": aes_key_for_api,
              "cs": filesize, "ps": rawsize, "fn": path.name, "md5": rawfilemd5}
    media_item = item_builder(**kwargs)

    msg = {
        "from_user_id": "", "to_user_id": to_user_id,
        "client_id": str(int(time.time() * 1000)),
        "message_type": MSG_TYPE_BOT, "message_state": MSG_STATE_FINISH,
        "item_list": [media_item],
    }
    if context_token:
        msg["context_token"] = context_token

    try:
        await _ilink_post("ilink/bot/sendmessage", {"msg": msg}, token, timeout_ms=15_000)
        # 先发文件再发说明文字
        if caption:
            await send_text(token, to_user_id, caption, context_token, bot_id=bot_id)
        root_log.info("[send_file] ✓ %s → %s", path.name, to_user_id[:20])
        return True
    except Exception as e:
        root_log.warning("[send_file] sendmessage 失败: %s", e)
        return False

_STATIC_DOWNLOADS = "/home/ubuntu/weclaw-1/app/static/downloads"
_STATIC_URL = "https://weclaw.pangoozn.com/static/downloads"


def _convert_media_to_http(text: str) -> str:
    """将 agent 回复中的 MEDIA:/path 转为 HTTP 链接（QQ 不支持 MEDIA 协议）。"""
    import re, os, shutil

    def _replace_media(m):
        raw_path = m.group(1).strip().strip("`\"'")
        if not os.path.exists(raw_path):
            alt = os.path.join(_STATIC_DOWNLOADS, raw_path)
            if os.path.exists(alt):
                raw_path = alt
            else:
                return m.group(0)
        fname = os.path.basename(raw_path)
        dst = os.path.join(_STATIC_DOWNLOADS, fname)
        try:
            if not os.path.exists(dst):
                os.makedirs(_STATIC_DOWNLOADS, exist_ok=True)
                shutil.copy2(raw_path, dst)
        except Exception:
            return m.group(0)
        return f"📎 {_STATIC_URL}/{fname}"

    _MEDIA_RE = re.compile(r'MEDIA:\s*([^\s,;)]+)')
    return _MEDIA_RE.sub(_replace_media, text)


async def send_text(token: str, to_user_id: str, text: str,
                    context_token: str = "",
                    session: Optional[aiohttp.ClientSession] = None,
                    bot_id: str = ""):
    """发送文字消息（自动将 MEDIA: 路径转为 HTTP 链接）。遇到 session 过期(-14)时去掉 token 重试一次。
    bot_id 传了就可以从 TokenStore 自动补 context_token。"""
    # 转换 MEDIA: 路径为 HTTP 链接（QQ 通道不支持 MEDIA 文件推送）
    if "MEDIA:" in text:
        text = _convert_media_to_http(text)
    msg = {
        "from_user_id": "", "to_user_id": to_user_id,
        "client_id": str(int(time.time() * 1000)),
        "message_type": MSG_TYPE_BOT, "message_state": MSG_STATE_FINISH,
        "item_list": [{"type": ITEM_TEXT, "text_item": {"text": text}}],
    }

    # 如果没传 context_token 但给了 bot_id，从 TokenStore 查
    ctx = context_token
    if not ctx and bot_id:
        ctx = _token_store.get(bot_id, to_user_id)

    if ctx:
        msg["context_token"] = ctx

    for attempt in range(2):  # 最多 2 次（第一次带 token，第二次不带）
        try:
            resp = await _ilink_post("ilink/bot/sendmessage", {"msg": msg}, token,
                                     timeout_ms=10_000, session=session)
            ret = resp.get("ret", 0)
            errcode = resp.get("errcode", 0)
            root_log.info("[send] → %s ret=%s errcode=%s", to_user_id[:20], ret, errcode)

            # Session 过期，去 token 重试
            if (ret == SESSION_EXPIRED_ERRCODE or errcode == SESSION_EXPIRED_ERRCODE) \
                    and attempt == 0 and ctx:
                root_log.warning("[send] session expired for %s, retrying without context_token", to_user_id[:20])
                msg.pop("context_token", None)
                ctx = ""
                if bot_id:
                    _token_store.remove(bot_id, to_user_id)
                await asyncio.sleep(0.5)
                continue

            # 第二次重试后还是 session 过期 → 真实失败
            if attempt == 1 and errcode == SESSION_EXPIRED_ERRCODE:
                root_log.warning("[send] 用户 %s session 过期，消息未送达", to_user_id[:20])
                return False

            # 其他错误 log 不退
            if ret not in (0, None) or errcode not in (0, None):
                root_log.warning("[send] err ret=%s errcode=%s", ret, errcode)
            return ret in (0, None)
        except Exception as e:
            root_log.warning("[send] 失败(attempt %d): %s", attempt + 1, e)
            if attempt == 0:
                await asyncio.sleep(1)
                continue
            return False
    return False

async def send_typing(token: str, to_user_id: str, typing_ticket: str,
                      status: int, session: Optional[aiohttp.ClientSession] = None):
    try:
        await _ilink_post("ilink/bot/sendtyping", {
            "ilink_user_id": to_user_id,
            "typing_ticket": typing_ticket,
            "status": status,
        }, token, timeout_ms=10_000, session=session)
    except Exception:
        pass

# ══════════════════════════════════════════════
# DB
# ══════════════════════════════════════════════

async def load_bots() -> list[tuple[str, str]]:
    try:
        import asyncpg
        conn = await asyncpg.connect(DB_DSN)
        try:
            rows = await conn.fetch(
                "SELECT bot_id, bot_token FROM bot_accounts "
                "WHERE is_active = true AND bot_token IS NOT NULL AND bot_token != ''"
            )
            return [(r["bot_id"], r["bot_token"]) for r in rows]
        finally:
            await conn.close()
    except Exception as e:
        root_log.error("DB 查询失败: %s", e)
        return []

# ══════════════════════════════════════════════
# 工具函数
# ══════════════════════════════════════════════

async def _wait(seconds: float, shutdown: asyncio.Event):
    try:
        await asyncio.wait_for(asyncio.get_event_loop().create_future(), timeout=seconds)
    except asyncio.TimeoutError:
        pass
    except asyncio.CancelledError:
        raise

# ══════════════════════════════════════════════
# Bot Session
# ══════════════════════════════════════════════

@dataclass
class BotSession:
    bot_id: str
    token: str
    log: logging.Logger
    sync_buf: str = ""
    last_activity: float = 0.0
    consecutive_failures: int = 0
    user_typing_tickets: dict = None  # user_id → typing_ticket
    welcome_sent: set = None  # set of user_ids that got welcome
    pending_quota_reminder: dict = None  # user_id → "剩余X次"

    def __post_init__(self):
        self.user_typing_tickets = self.user_typing_tickets or {}
        self.welcome_sent = self.welcome_sent or set()
        self.pending_quota_reminder = self.pending_quota_reminder or {}

# ══════════════════════════════════════════════
# Agent Connector HTTP 调用
# ══════════════════════════════════════════════

async def forward_to_agent(bot_id: str, from_user: str, text: str,
                           msg_id: str, context_token: str) -> Optional[str]:
    """POST 消息到 Agent Connector，等回复"""
    try:
        async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=AGENT_TIMEOUT)) as s:
            async with s.post(AGENT_URL, json={
                "bot_id": bot_id,
                "from_user": from_user,
                "text": text,
                "msg_id": msg_id,
                "context_token": context_token,
            }) as r:
                if r.status == 200:
                    data = await r.json()
                    return data.get("reply")
    except asyncio.TimeoutError:
        root_log.warning("[agent] 超时(120s)")
    except Exception as e:
        root_log.warning("[agent] 失败: %s", e)
    return None

# ══════════════════════════════════════════════
# 主轮询循环
# ══════════════════════════════════════════════

async def run_bot(bot_id: str, token: str, shutdown: asyncio.Event):
    log = _make_logger(f"bot:{bot_id[:14]}")
    bot = BotSession(bot_id=bot_id, token=token, log=log)
    register_bot(bot_id, bot)
    log.info("🤖 启动")

    # 预热
    has_old = bool(load_sync_buf(bot_id))
    if has_old:
        await notify_stop(token)
        await asyncio.sleep(2)
    else:
        log.info("  新 Bot，跳过 notifyStop")

    try:
        # ... session 循环 ...
        async with aiohttp.ClientSession() as session:
            while not shutdown.is_set():
                try:
                    await _run_session(bot, session, shutdown)
                except asyncio.CancelledError:
                    break
                except Exception as e:
                    log.error("❌ session 崩溃: %s", e)
                    bot.consecutive_failures += 1
                    if bot.consecutive_failures >= 3:
                        log.error("连续 3 次崩溃，暂停 10min")
                        await _wait(600, shutdown)
                        bot.consecutive_failures = 0

        await notify_stop(token)
    finally:
        unregister_bot(bot_id)
    log.info("🤖 已停止")

async def _run_session(bot: BotSession, session: aiohttp.ClientSession,
                       shutdown: asyncio.Event):
    log = bot.log
    bot.last_activity = time.time()

    # 恢复 sync_buf
    bot.sync_buf = load_sync_buf(bot.bot_id)
    if bot.sync_buf:
        log.info("  恢复 sync_buf (%d bytes)", len(bot.sync_buf))

    # notifyStart
    if not await notify_start(bot.token):
        log.warning("  notifyStart 失败，重试...")
        await _wait(5, shutdown)
        if not await notify_start(bot.token):
            log.error("  notifyStart 连续失败，跳过")
            return
    log.info("✅ notifyStart ✓")

    next_timeout = LONG_POLL_TIMEOUT_MS

    while not shutdown.is_set():
        # 暂停检查
        if is_paused(bot.bot_id):
            remaining = pause_remaining_ms(bot.bot_id)
            log.info("⏳ 暂停中，剩余 %dmin", remaining // 60_000)
            await _wait(min(remaining / 1000, 600), shutdown)
            clear_pause(bot.bot_id)
            bot.consecutive_failures = 0
            await notify_stop(bot.token)
            await asyncio.sleep(2)
            return  # 外层重建 session

        try:
            resp = await get_updates(bot.token, bot.sync_buf, next_timeout, session)
            ret = resp.get("ret", 0)
            errcode = resp.get("errcode", 0)
            is_error = (ret != 0) or (errcode != 0)

            if is_error:
                if errcode == SESSION_EXPIRED_ERRCODE or ret == SESSION_EXPIRED_ERRCODE:
                    log.error("🚫 session 过期，暂停 10min")
                    pause_bot(bot.bot_id)
                    await _wait(10, shutdown)
                    bot.consecutive_failures = 0
                    return
                bot.consecutive_failures += 1
                log.warning("getUpdates 错误 ret=%s err=%s (%d/3)", ret, errcode, bot.consecutive_failures)
                if bot.consecutive_failures >= 3:
                    await _wait(30, shutdown)
                    bot.consecutive_failures = 0
                else:
                    await _wait(2, shutdown)
                continue

            # 成功
            bot.consecutive_failures = 0
            bot.last_activity = time.time()

            if resp.get("longpolling_timeout_ms"):
                next_timeout = resp["longpolling_timeout_ms"]

            # sync_buf 持久化
            new_buf = resp.get("get_updates_buf", "")
            if new_buf and new_buf != bot.sync_buf:
                save_sync_buf(bot.bot_id, new_buf)
                bot.sync_buf = new_buf

            # 处理消息
            for msg in resp.get("msgs", []):
                from_user = msg.get("from_user_id", "")
                if not from_user:
                    continue

                ctx = msg.get("context_token", "")
                # 持久化 context_token（用于后续主动推送）
                if ctx:
                    _token_store.set(bot.bot_id, from_user, ctx)
                text = _extract_text(msg, from_user)

                # 文件/图片/视频/语音消息 → 异步更新磁盘配额
                if text and (text.startswith("[文件]") or text.startswith("[图片]") or text.startswith("[视频]") or text.startswith("[语音]")):
                    asyncio.create_task(_update_disk_quota(from_user))

                if not text:
                    continue

                log.info("📩 %s: %s", from_user[:20], text[:60])

                # ── 1. 暗号匹配（秒回，不经过 Agent）──
                code_match = match_code_phrase(text)
                if code_match:
                    reply = code_match["reply"]
                    log.info("🔐 暗号: %s → %s", text[:12], reply)
                    await send_text(bot.token, from_user, reply, ctx, session)
                    continue

                # ── 1.3 项目口令匹配（秒回加入项目）──
                import re as _re
                _invite_match = _re.match(r'^(?:口令|项目)[-：:](.+)$', text.strip())
                if _invite_match:
                    code_suffix = _invite_match.group(1)
                    code_key = f"项目-{code_suffix}"
                    uid = from_user.split("@")[0]
                    # 查昵称
                    _invite_nick = uid[:8]
                    try:
                        import asyncpg
                        _invite_db = await asyncpg.connect(DB_DSN)
                        try:
                            _invite_row = await _invite_db.fetchrow(
                                "SELECT nickname FROM channel_bindings WHERE channel_user_id LIKE $1 AND nickname IS NOT NULL AND nickname != '' LIMIT 1",
                                uid + "%"
                            )
                            if _invite_row and _invite_row["nickname"]:
                                _invite_nick = _invite_row["nickname"]
                        finally:
                            await _invite_db.close()
                    except Exception:
                        pass
                    try:
                        import aiohttp
                        async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=10)) as _s:
                            async with _s.post(
                                "https://hai.pangoozn.com/api/projects/join",
                                json={"code": code_key, "nickname": _invite_nick, "channel": "139bot"},
                                headers={"Content-Type": "application/json"}
                            ) as _r:
                                if _r.status == 200:
                                    _jd = await _r.json()
                                    reply = _jd.get("message", f"✅ 已加入项目")
                                else:
                                    _jt = await _r.text()
                                    reply = f"❌ 加入失败（{_r.status}）"
                    except Exception as _ie:
                        root_log.error("项目口令异常: %s", _ie)
                        reply = "❌ 项目服务不可达"
                    await send_text(bot.token, from_user, reply, ctx, session)
                    continue

                # ── 1.4 项目命令查询 + 管理（秒回，不经过 Agent）──
                _project_cmd = text.strip()
                _project_lower = _project_cmd.lower()

                # 所有项目命令都需要查询用户昵称
                uid = from_user.split("@")[0]
                _nick = uid[:8]
                try:
                    import asyncpg
                    _nick_db = await asyncpg.connect(DB_DSN)
                    try:
                        _nick_row = await _nick_db.fetchrow(
                            "SELECT nickname FROM channel_bindings WHERE channel_user_id LIKE $1 AND nickname IS NOT NULL AND nickname != '' LIMIT 1",
                            uid + "%"
                        )
                        if _nick_row and _nick_row["nickname"]:
                            _nick = _nick_row["nickname"]
                    finally:
                        await _nick_db.close()
                except Exception:
                    pass

                # --- 项目列表（支持"项目列表"和"项目清单"）---
                if _project_lower in ("项目列表", "项目清单"):
                    try:
                        import aiohttp, urllib.parse
                        async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=15)) as _s:
                            async with _s.get(f"https://hai.pangoozn.com/api/projects/list?nickname={urllib.parse.quote(_nick)}") as _r:
                                if _r.status == 200:
                                    _jd = await _r.json()
                                    projects = _jd.get("projects", [])
                                    if projects:
                                        reply = "📁 你的项目：\n" + "\n".join(f"  {p}" for p in projects)
                                    else:
                                        reply = "📁 你还没有加入任何项目。需要口令才能加入"
                                else:
                                    _jt = await _r.text()
                                    reply = f"❌ 查询失败（{_r.status}）"
                    except Exception as _pe:
                        root_log.error("项目列表异常: %s", _pe)
                        reply = "❌ 项目服务不可达"
                    await send_text(bot.token, from_user, reply, ctx, session)
                    continue

                # --- 项目成员 ---
                if _project_lower == "项目成员":
                    try:
                        import aiohttp, urllib.parse
                        async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=15)) as _s:
                            async with _s.get(f"https://hai.pangoozn.com/api/projects/members?nickname={urllib.parse.quote(_nick)}") as _r:
                                if _r.status == 200:
                                    _jd = await _r.json()
                                    members = _jd.get("members", [])
                                    proj = _jd.get("project", "未知")
                                    if members:
                                        lines = [f"📂 {proj} — {len(members)} 位成员："]
                                        for m in members:
                                            role_str = "👑" if m.get("role") == "admin" else "👤"
                                            nick = m.get("nickname", "?")
                                            ch = m.get("channel", "")
                                            lines.append(f"  {role_str} {nick} ({ch})")
                                        reply = "\n".join(lines)
                                    else:
                                        reply = f"📂 {proj} — 暂无成员"
                                else:
                                    _jt = await _r.text()
                                    reply = f"❌ 查询失败（{_r.status}）"
                    except Exception as _pe:
                        root_log.error("项目成员异常: %s", _pe)
                        reply = "❌ 项目服务不可达"
                    await send_text(bot.token, from_user, reply, ctx, session)
                    continue

                # --- 项目文件（支持子目录：项目文件 目录名 / 文件列表 目录名）---
                _file_cmd = re.match(r'^(?:项目文件|文件列表)(?:\s+(.+))?$', _project_cmd)
                if _file_cmd:
                    sub_dir = _file_cmd.group(1) or ""
                    try:
                        import aiohttp, urllib.parse
                        async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=15)) as _s:
                            _url = f"https://hai.pangoozn.com/api/projects/files?nickname={urllib.parse.quote(_nick)}&dir={urllib.parse.quote(sub_dir)}"
                            async with _s.get(_url) as _r:
                                if _r.status == 200:
                                    _jd = await _r.json()
                                    entries = _jd.get("entries", [])
                                    proj = _jd.get("project", "未知")
                                    cur_dir = _jd.get("dir", "")
                                    if entries:
                                        header = f"📂 {proj}"
                                        if cur_dir:
                                            header += f"/{cur_dir}"
                                        header += f"（{len(entries)} 项）"
                                        lines = [header]
                                        for e in entries:
                                            icon = "📁" if e.get("type") == "dir" else "📄"
                                            size = e.get("size", 0)
                                            size_str = f"{size/1024:.0f}KB" if size > 1024 else f"{size}B"
                                            lines.append(f"  {icon} {e['name']} ({size_str})")
                                        reply = "\n".join(lines)
                                    else:
                                        reply = f"📂 {proj} — 空目录"
                                else:
                                    _jt = await _r.text()
                                    reply = f"❌ 查询失败（{_r.status}）: {_jt[:100]}"
                    except Exception as _pe:
                        root_log.error("项目文件异常: %s", _pe)
                        reply = "❌ 项目服务不可达"
                    await send_text(bot.token, from_user, reply, ctx, session)
                    continue

                # --- 进入项目（进入 XX 或 进入项目 XX）---
                _enter_cmd = re.match(r'^(?:进入(?:项目)?)\s*(.+)?$', _project_cmd)
                if _enter_cmd:
                    _proj_name = _enter_cmd.group(1)
                    if not _proj_name:
                        reply = "📝 格式：进入 [项目名]  或  进入项目 [项目名]"
                    else:
                        _proj_name = _proj_name.strip()
                        try:
                            import aiohttp
                            async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=10)) as _s:
                                async with _s.post(
                                    "https://hai.pangoozn.com/api/projects/enter",
                                    json={"name": _proj_name, "nickname": _nick, "channel": "139bot"},
                                    headers={"Content-Type": "application/json"}
                                ) as _r:
                                    if _r.status == 200:
                                        reply = f"✅ 已进入「{_proj_name}」"
                                        # 进入成功后自动拉文件列表
                                        import urllib.parse
                                        try:
                                            async with _s.get(f"https://hai.pangoozn.com/api/projects/files?nickname={urllib.parse.quote(_nick)}") as _fr:
                                                if _fr.status == 200:
                                                    _fd = await _fr.json()
                                                    entries = _fd.get("entries", [])
                                                    if entries:
                                                        lines = [f"\n📂 {_proj_name} 文件："]
                                                        for e in entries:
                                                            icon = "📁" if e.get("type") == "dir" else "📄"
                                                            sz = e.get("size", 0)
                                                            szs = f"{sz/1024:.0f}KB" if sz > 1024 else f"{sz}B"
                                                            lines.append(f"  {icon} {e['name']} ({szs})")
                                                        reply += "\n" + "\n".join(lines)
                                                    else:
                                                        reply += "\n📂 项目暂无文件"
                                        except Exception:
                                            pass
                                    elif _r.status == 403:
                                        _jt = await _r.json()
                                        reply = _jt.get("detail", "❌ 无权进入该项目，请先通过口令加入")
                                    else:
                                        _jt = await _r.text()
                                        reply = f"❌ 进入失败（{_r.status}）"
                        except Exception as _ee:
                            root_log.error("进入项目异常: %s", _ee)
                            reply = "❌ 项目服务不可达"
                    await send_text(bot.token, from_user, reply, ctx, session)
                    continue

                # --- 看看文件（模糊搜索项目文件并读取内容）---
                _lookup_cmd = re.match(r'^(?:看看|查看|打开)\s+(.+)$', _project_cmd)
                if _lookup_cmd:
                    _query = _lookup_cmd.group(1).strip()
                    try:
                        import aiohttp, urllib.parse
                        async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=20)) as _s:
                            _search_body = json.dumps({"keyword": _query, "nickname": _nick}).encode()
                            async with _s.post("https://hai.pangoozn.com/api/projects/search", json={"keyword": _query, "nickname": _nick}) as _sr:
                                if _sr.status == 200:
                                    _sd = await _sr.json()
                                    _results = _sd.get("results", [])
                                    if not _results:
                                        reply = "🔍 未找到包含「{q}」的文件".format(q=_query)
                                    elif len(_results) == 1:
                                        _match = _results[0]
                                        _proj = _match["project"]
                                        _fname = _match["file"]
                                        _read_url = "https://hai.pangoozn.com/api/projects/read?project={p}&file={f}&nickname={n}".format(
                                            p=urllib.parse.quote(_proj), f=urllib.parse.quote(_fname), n=urllib.parse.quote(_nick))
                                        async with _s.get(_read_url) as _rr:
                                            if _rr.status == 200:
                                                _rd = await _rr.json()
                                                _content = _rd.get("content", "")
                                                if len(_content) > 800:
                                                    _content = _content[:800] + "\n\n...(截断)"
                                                _sep = "=" * 40
                                                reply = "📄 {p}/{f}\n{sep}\n{c}".format(p=_proj, f=_fname, sep=_sep, c=_content)
                                            else:
                                                reply = "❌ 读取失败"
                                    else:
                                        # 多个匹配，读第一个（最相关）
                                        _match = _results[0]
                                        _proj = _match["project"]
                                        _fname = _match["file"]
                                        _read_url = "https://hai.pangoozn.com/api/projects/read?project={p}&file={f}&nickname={n}".format(
                                            p=urllib.parse.quote(_proj), f=urllib.parse.quote(_fname), n=urllib.parse.quote(_nick))
                                        async with _s.get(_read_url) as _rr:
                                            if _rr.status == 200:
                                                _rd = await _rr.json()
                                                _content = _rd.get("content", "")
                                                if len(_content) > 800:
                                                    _content = _content[:800] + "\n\n...(截断)"
                                                _sep = "=" * 40
                                                reply = "📄 {p}/{f}\n{sep}\n{c}".format(p=_proj, f=_fname, sep=_sep, c=_content)
                                            else:
                                                reply = "❌ 读取失败"
                                else:
                                    reply = "❌ 搜索失败"
                    except Exception as _pe:
                        root_log.error("看看文件异常: %s", _pe)
                        reply = "❌ 项目服务不可达"
                    await send_text(bot.token, from_user, reply, ctx, session)
                    continue

                # --- 生成口令（管理员）---
                if _project_lower == "生成口令":
                    try:
                        import aiohttp
                        async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=10)) as _s:
                            async with _s.post(
                                "https://hai.pangoozn.com/api/projects/invite",
                                json={"nickname": _nick, "admin": _nick},
                                headers={"Content-Type": "application/json"}
                            ) as _r:
                                if _r.status == 200:
                                    _jd = await _r.json()
                                    reply = f"🔑 邀请口令：{_jd.get('invite_code', '?')}"
                                elif _r.status == 403:
                                    _jt = await _r.json()
                                    reply = _jt.get("detail", "❌ 仅项目管理员可生成口令")
                                else:
                                    _jt = await _r.text()
                                    reply = f"❌ 生成失败（{_r.status}）"
                    except Exception as _ge:
                        root_log.error("生成口令异常: %s", _ge)
                        reply = "❌ 项目服务不可达"
                    await send_text(bot.token, from_user, reply, ctx, session)
                    continue

                # --- 踢出成员（管理员）---
                _kick_cmd = re.match(r'^踢出\s+(.+)$', _project_cmd)
                if _kick_cmd:
                    _target = _kick_cmd.group(1).strip()
                    try:
                        import aiohttp
                        async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=10)) as _s:
                            async with _s.post(
                                "https://hai.pangoozn.com/api/projects/members/remove",
                                json={"nickname": _target, "operator": _nick},
                                headers={"Content-Type": "application/json"}
                            ) as _r:
                                if _r.status == 200:
                                    reply = f"✅ 已移除成员「{_target}」"
                                elif _r.status == 403:
                                    _jt = await _r.json()
                                    reply = _jt.get("detail", "❌ 仅项目管理员可踢出成员")
                                else:
                                    _jt = await _r.text()
                                    reply = f"❌ 移除失败（{_r.status}）"
                    except Exception as _ke:
                        root_log.error("踢出成员异常: %s", _ke)
                        reply = "❌ 项目服务不可达"
                    await send_text(bot.token, from_user, reply, ctx, session)
                    continue

                # --- 创建项目 ---
                _create_cmd = re.match(r'^创建项目\s+(.+)$', _project_cmd)
                if _create_cmd:
                    _proj_name = _create_cmd.group(1).strip()
                    try:
                        import aiohttp
                        async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=10)) as _s:
                            async with _s.post(
                                "https://hai.pangoozn.com/api/projects/create",
                                json={"name": _proj_name, "nickname": _nick, "channel": "139bot"},
                                headers={"Content-Type": "application/json"}
                            ) as _r:
                                if _r.status == 200:
                                    _jd = await _r.json()
                                    reply = f"✅ {_jd.get('message', '项目已创建')}"
                                elif _r.status == 409:
                                    _jt = await _r.json()
                                    reply = f"❌ {_jt.get('detail', '项目已存在')}"
                                else:
                                    _jt = await _r.text()
                                    reply = f"❌ 创建失败（{_r.status}）"
                    except Exception as _ce:
                        root_log.error("创建项目异常: %s", _ce)
                        reply = "❌ 项目服务不可达"
                    await send_text(bot.token, from_user, reply, ctx, session)
                    continue

                # --- 设为管理员（仅超级管理员可操作）---
                _promote_cmd = re.match(r'^设为管理员\s+(.+)$', _project_cmd)
                if _promote_cmd:
                    _target = _promote_cmd.group(1).strip()
                    try:
                        import aiohttp
                        async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=10)) as _s:
                            async with _s.post(
                                "https://hai.pangoozn.com/api/projects/members/promote",
                                json={"nickname": _target, "operator": _nick},
                                headers={"Content-Type": "application/json"}
                            ) as _r:
                                if _r.status == 200:
                                    _jd = await _r.json()
                                    reply = f"✅ {_jd.get('message', f'{_target} 已提升为管理员')}"
                                elif _r.status == 403:
                                    _jt = await _r.json()
                                    reply = _jt.get("detail", "❌ 仅超级管理员可执行此操作")
                                else:
                                    _jt = await _r.text()
                                    reply = f"❌ 提升失败（{_r.status}）"
                    except Exception as _pe:
                        root_log.error("设为管理员异常: %s", _pe)
                        reply = "❌ 项目服务不可达"
                    await send_text(bot.token, from_user, reply, ctx, session)
                    continue

                # ── @weclaw 召唤主脑（提交+回调）──
                if text.strip().startswith("@weclaw"):
                    main_text = text.replace("@weclaw", "", 1).strip()
                    import aiohttp
                    try:
                        async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=10)) as s:
                            async with s.post(
                                "https://hai.pangoozn.com/api/weclaw-proxy",
                                json={
                                    "bot_id": bot.bot_id,
                                    "user_id": from_user,
                                    "content": main_text or "召唤主脑",
                                    "msg_id": msg.get("msg_id", ""),
                                    "callback_url": "https://ai.pangoozn.com/api/bot/weclaw-callback",
                                    "callback_bot_id": bot.bot_id,
                                    "callback_to_user": from_user,
                                    "callback_ctx": ctx,
                                },
                                headers={"Content-Type": "application/json"}
                            ) as r:
                                if r.status == 200:
                                    data = await r.json()
                                    if data.get("success"):
                                        await send_text(bot.token, from_user, "🧠 已召唤主脑，稍后回复~", ctx, session)
                                    else:
                                        await send_text(bot.token, from_user, "🧠 主脑提交失败", ctx, session)
                                else:
                                    await send_text(bot.token, from_user, f"🧠 主脑不可达（{r.status}）", ctx, session)
                    except Exception as e:
                        root_log.error("@weclaw 异常: %s", e)
                        await send_text(bot.token, from_user, "🧠 主脑召唤失败", ctx, session)
                    # @weclaw 消息只发给主脑，不落本地 Agent
                    continue
                    # 回调会通过 /api/bot/weclaw-callback 推回结果

                # ── 1.4 会员指令（查询/续费，不经过 Agent）──
                _cmd_text = text.strip().lower()
                if _cmd_text in ("会员", "我的套餐", "套餐", "查会员", "查套餐", "vip", "我的会员"):
                    uid = from_user.split("@")[0]
                    try:
                        import asyncpg
                        _qdb = await asyncpg.connect(DB_DSN)
                        try:
                            # 先查 channel_bindings 转成服务号 openid
                            svc_openid = uid
                            cb = await _qdb.fetchrow(
                                "SELECT openid FROM channel_bindings WHERE channel_user_id LIKE $1 LIMIT 1",
                                uid + "%"
                            )
                            if cb and cb["openid"]:
                                svc_openid = cb["openid"]
                            root_log.info("会员查询: uid=%s → svc_openid=%s", uid, svc_openid)
                            row = await _qdb.fetchrow(
                                "SELECT s.status, s.expires_at, p.name as plan_name, s.started_at, s.xiake_points, s.total_points_consumed, s.disk_used_bytes, s.disk_quota_bytes "
                                "FROM subscribers s LEFT JOIN plans p ON s.plan_id = p.id "
                                "WHERE s.openid = $1", svc_openid
                            )
                            if row and str(row["status"]).upper() in ("ACTIVE", "TRIAL") and row["expires_at"] and row["expires_at"] >= date.today():
                                remain_days = (row["expires_at"] - date.today()).days
                                root_log.info("会员查询: ✅ 找到会员 openid=%s status=%s exp=%s remain=%d", svc_openid, row["status"], row["expires_at"], remain_days)
                                pts = row["xiake_points"] or 0
                                consumed = row["total_points_consumed"] or 0
                                disk_mb = (row["disk_used_bytes"] or 0) // (1024*1024)
                                disk_total = (row["disk_quota_bytes"] or 2147483648) // (1024*1024)
                                reply_text = (
                                    f"🦞 享客虾会员 · {row['plan_name'] or '已开通'}\n\n"
                                    f"📅 到期：{row['expires_at']}（剩余 {remain_days} 天）\n"
                                    f"🦞 虾点：{pts} 点\n"
                                    f"📊 累计消耗：{consumed} 点\n"
                                    f"💾 磁盘：{disk_mb}MB / {disk_total}MB\n"
                                    f"📆 开通：{row['started_at']}\n\n"
                                    f"👉 续期会员：https://ai.pangoozn.com/subscribe"
                                )
                            else:
                                reply_text = (
                                    f"🦞 你当前是免费用户\n\n"
                                    f"你可以免费浏览享客虾的全部作品案例，\n"
                                    f"开通会员即可解锁创作功能 ✨\n\n"
                                    f"👉 https://ai.pangoozn.com/subscribe"
                                )
                        finally:
                            await _qdb.close()
                    except Exception as _qe:
                        root_log.warning("会员查询异常: %s", _qe)
                        reply_text = "⚠️ 查询失败，请稍后再试"
                    
                    await send_text(bot.token, from_user, reply_text, ctx, session)
                    continue

                if _cmd_text in ("续费", "开通", "订阅", "购买"):
                    reply_text = (
                        f"🦞 开通享客虾会员\n\n"
                        f"• 月卡 ¥99/月（10000虾点）\n"
                        f"• 年卡 ¥999/年（120000虾点，省¥189）\n\n"
                        f"👇 点此开通\n"
                        f"https://ai.pangoozn.com/subscribe\n\n"
                        f"回复「会员」查看当前状态"
                    )
                    await send_text(bot.token, from_user, reply_text, ctx, session)
                    continue

                # ── 1.5 欢迎消息（首次自动 + hi/你好 触发）──
                _greet = _project_cmd.strip().lower()
                _is_greeting = _greet in ("hi", "hello", "你好", "嗨", "在吗", "在不在", "help", "帮助")
                # 检查是否已欢迎过（DB 持久化）
                _already_welcomed = False
                try:
                    import asyncpg
                    _wuid = from_user.split("@")[0]
                    _wdb = await asyncpg.connect(DB_DSN)
                    try:
                        _wrow = await _wdb.fetchrow(
                            "SELECT welcomed FROM channel_bindings WHERE channel_user_id LIKE $1 AND welcomed = true LIMIT 1",
                            _wuid + "%"
                        )
                        if _wrow:
                            _already_welcomed = True
                    finally:
                        await _wdb.close()
                except Exception:
                    pass

                # 条件：首次任何消息 或 hi类命令
                if (not _already_welcomed) or _is_greeting:
                    _uid = from_user.split("@")[0]
                    _nick = _uid[:8]
                    _member_line = ""
                    try:
                        from datetime import date
                        _wdb = await asyncpg.connect(DB_DSN)
                        try:
                            _wrow = await _wdb.fetchrow(
                                "SELECT nickname FROM channel_bindings WHERE channel_user_id LIKE $1 AND nickname IS NOT NULL AND nickname != '' LIMIT 1",
                                _uid + "%"
                            )
                            if _wrow and _wrow["nickname"]:
                                _nick = _wrow["nickname"]
                            cb2 = await _wdb.fetchrow(
                                "SELECT openid FROM channel_bindings WHERE channel_user_id LIKE $1 LIMIT 1",
                                _uid + "%"
                            )
                            svc_openid = cb2["openid"] if (cb2 and cb2["openid"]) else _uid
                            _mrow = await _wdb.fetchrow(
                                "SELECT s.status, s.expires_at, p.name as plan_name "
                                "FROM subscribers s LEFT JOIN plans p ON s.plan_id = p.id "
                                "WHERE s.openid = $1", svc_openid
                            )
                            if _mrow and str(_mrow["status"]).upper() in ("ACTIVE", "TRIAL") and _mrow["expires_at"] and _mrow["expires_at"] >= date.today():
                                _remain = (_mrow["expires_at"] - date.today()).days
                                _member_line = (
                                    f"\n📅 会员到期：{_mrow['expires_at']}（剩余 {_remain} 天）\n"
                                    f"👉 续费：https://ai.pangoozn.com/subscribe"
                                )
                            else:
                                _member_line = (
                                    f"\n💎 开通会员 → https://ai.pangoozn.com/subscribe"
                                )
                            # 标记已欢迎（仅首次）
                            if not _already_welcomed:
                                await _wdb.execute(
                                    "UPDATE channel_bindings SET welcomed = true WHERE channel_user_id LIKE $1",
                                    _uid + "%"
                                )
                        finally:
                            await _wdb.close()
                    except Exception:
                        pass
                    reply = (
                        f"✨ 欢迎你，{_nick}！\n\n"
                        f"🦞 享客虾 Bot 已就绪 ✅"
                        f"{_member_line}\n\n"
                        f"有需要随时招呼，试试发「帮助」了解我能做什么"
                    )
                    await send_text(bot.token, from_user, reply, ctx, session)
                    # 不 continue！首次消息和 hi 命令都继续走 AI 处理

                # ── 1.6 会员/配额检查 ──
                uid = from_user.split("@")[0]
                _q_member = False
                _points_note = ""
                svc_openid = uid  # 补充：待查 channel_bindings 获取服务号 openid
                try:
                    import asyncpg
                    _q_db = await asyncpg.connect(DB_DSN)
                    try:
                        # 先查 channel_bindings 获取服务号 openid
                        cb_row = await _q_db.fetchrow(
                            "SELECT openid FROM channel_bindings WHERE channel_user_id LIKE $1 LIMIT 1",
                            uid + "%"
                        )
                        if cb_row and cb_row["openid"]:
                            svc_openid = cb_row["openid"]

                        # 用服务号 openid 查 subscriber（优先），没有则用 iLink ID
                        row = await _q_db.fetchrow(
                            "SELECT status, expires_at, xiake_points FROM subscribers WHERE openid = $1",
                            svc_openid
                        )
                        if not row:
                            row = await _q_db.fetchrow(
                                "SELECT status, expires_at, xiake_points FROM subscribers WHERE openid = $1",
                                uid
                            )
                        if row:
                            st = str(row["status"]).upper()
                            exp = row["expires_at"]
                            _q_member = st in ("ACTIVE", "TRIAL") and exp and exp >= date.today()

                        if _q_member:
                            # ── 会员走虾点 ──
                            pts = row["xiake_points"] if row else 0
                            if pts < 1:
                                await send_text(bot.token, from_user,
                                    f"🦞 虾点已用完（剩余 {pts} 点）~\n\n"
                                    f"充值入口 👉 https://ai.pangoozn.com/subscribe\n"
                                    f"续费会员畅享 10000 点/月无限对话 ✨",
                                    ctx, session)
                                await _q_db.close()
                                continue

                            # 扣 1 虾点
                            new_pts = pts - 1
                            await _q_db.execute(
                                "UPDATE subscribers SET xiake_points = $1, "
                                "total_points_consumed = total_points_consumed + 1 "
                                "WHERE openid = $2 AND xiake_points = $3",
                                new_pts, svc_openid, pts
                            )
                            # 写流水
                            sub_id = await _q_db.fetchval(
                                "SELECT id FROM subscribers WHERE openid = $1", svc_openid
                            )
                            if sub_id:
                                await _q_db.execute(
                                    "INSERT INTO points_transactions (subscriber_id, tx_type, amount, "
                                    "balance_after, description) VALUES ($1, 'consume', -1, $2, 'chat')",
                                    sub_id, new_pts
                                )

                            # 低点提醒
                            if new_pts <= 50:
                                bot.pending_quota_reminder[from_user] = f"🦞 虾点剩余 {new_pts} 点，点此充值 👉 https://ai.pangoozn.com/subscribe"
                            elif new_pts <= 300:
                                bot.pending_quota_reminder[from_user] = f"🦞 剩余 {new_pts} 点"

                        if not _q_member:
                            await send_text(bot.token, from_user,
                                "🦞 享客虾是付费服务，开通即可畅聊所有功能 ✨\n"
                                f"👉 https://ai.pangoozn.com/subscribe",
                                ctx, session)
                            await _q_db.close()
                            continue
                    finally:
                        await _q_db.close()
                except Exception as _qe:
                    log.warning("配额检查异常(放行): %s", _qe)

                # ── 2. 获取 typing ticket（每次刷新）──
                ticket = ""
                try:
                    cfg = await _ilink_get_config(bot.token, from_user, ctx or None, session)
                    ticket = str(cfg.get("typing_ticket") or "")
                    if ticket:
                        bot.user_typing_tickets[from_user] = ticket
                except Exception:
                    ticket = bot.user_typing_tickets.get(from_user, "")

                # ── 3. 发 typing + keep-alive ──
                typing_keepalive = None
                if ticket:
                    try:
                        await send_typing(bot.token, from_user, ticket, TYPING_START, session)
                        # keep-alive: 每 15s 发一次 TYPING_START，直到取消
                        async def _typing_keepalive():
                            try:
                                while not shutdown.is_set():
                                    await asyncio.sleep(15)
                                    if not ticket:
                                        break
                                    await send_typing(bot.token, from_user, ticket, TYPING_START, session)
                            except asyncio.CancelledError:
                                pass
                            except Exception:
                                pass
                        typing_keepalive = asyncio.create_task(_typing_keepalive())
                    except Exception as e:
                        log.warning("TYPING_START 失败: %s", e)

                # ── 4. 转发 Agent Connector ──
                try:
                    reply = await forward_to_agent(bot.bot_id, from_user, text,
                                                    msg.get("msg_id", ""), ctx)
                finally:
                    # ── 5. 停止 typing + 取消 keep-alive ──
                    if typing_keepalive:
                        typing_keepalive.cancel()
                    if ticket:
                        try:
                            await send_typing(bot.token, from_user, ticket, TYPING_STOP, session)
                        except Exception as e:
                            log.warning("TYPING_STOP 失败: %s", e)

                # ── 6. 发回复 ──
                if reply:
                    # 附上配额提示
                    qr = bot.pending_quota_reminder.pop(from_user, None)
                    if qr:
                        reply = reply + "\n\n" + qr

                    # 检测 MEDIA: 路径 → 文件推送（支持嵌入文中的多条 MEDIA）
                    media_paths = re.findall(r'MEDIA:(\S+?)(?:[\*\s）\)]|$)', reply)
                    if media_paths:
                        clean = reply
                        ok_count = 0
                        for mp in media_paths:
                            mp_clean = re.sub(r'[\*\）\)]+$', '', mp)
                            # 路径不存在时自动搜索已知目录
                            if not os.path.exists(mp_clean):
                                found = _search_file(mp_clean)
                                if found:
                                    mp_clean = found
                            ok = await _send_file(bot.token, from_user, mp_clean, "", ctx, bot_id=bot.bot_id)
                            if ok:
                                ok_count += 1
                            clean = clean.replace(f"MEDIA:{mp}", "").strip()
                            await asyncio.sleep(0.5)
                        # 清理 MEDIA 残留的 markdown 符号（**）和文件大小（（324KB））
                        clean = re.sub(r'\*+MEDIA:', '', clean)
                        clean = re.sub(r'\*+（[^）]*[KkMmGg][Bb]）', '', clean)
                        clean = re.sub(r'\*+', '', clean)
                        # 清理多余空行和括号残留
                        clean = re.sub(r'\n{3,}', '\n\n', clean)
                        clean = re.sub(r'[（(]\s*[)）]', '', clean)
                        if clean.strip() and ok_count > 0:
                            await send_text(bot.token, from_user, clean.strip(), ctx, session)
                        log.info("📤 [MEDIA] %d 个文件推送成功", ok_count)
                    else:
                        log.info("📤 %s", reply[:60])
                        await send_text(bot.token, from_user, reply, ctx, session)
                else:
                    log.warning("⚠️ Agent 无回复")

        except asyncio.CancelledError:
            raise
        except Exception as e:
            bot.consecutive_failures += 1
            log.error("轮询异常 (%d/3): %s", bot.consecutive_failures, e)
            if bot.consecutive_failures >= 3:
                await _wait(30, shutdown)
                bot.consecutive_failures = 0
            else:
                await _wait(2, shutdown)

# ══════════════════════════════════════════════
# 文件下载与内容提取
# ══════════════════════════════════════════════

import base64

FILE_DIR = os.path.join(os.path.dirname(__file__), "downloads")
os.makedirs(FILE_DIR, exist_ok=True)

# 文件类型 → 子目录名
_FILE_TYPE_DIRS = {
    "file": "files",
    "image": "images",
    "video": "video",
    "voice": "voice",
}

# 异步更新 DB disk_used_bytes（消息处理后调用）
async def _update_disk_quota(user_id: str) -> None:
    """扫描用户目录文件大小，更新 DB disk_used_bytes"""
    import asyncpg
    uid = user_id.split("@")[0]
    try:
        safe_uid = re.sub(r'[^\w\.\-@]', '_', user_id.split('/')[-1])[:64]
        user_dir = os.path.join(FILE_DIR, safe_uid)
        total = 0
        if os.path.isdir(user_dir):
            for root, dirs, files in os.walk(user_dir):
                for f in files:
                    try:
                        total += os.path.getsize(os.path.join(root, f))
                    except:
                        pass
        conn = await asyncpg.connect(DB_DSN)
        try:
            # 先通过 channel_bindings 转成服务号 openid
            svc_openid = uid
            cb = await conn.fetchrow(
                "SELECT openid FROM channel_bindings WHERE channel_user_id LIKE $1 LIMIT 1",
                uid + "%"
            )
            if cb and cb["openid"]:
                svc_openid = cb["openid"]
            await conn.execute(
                "UPDATE subscribers SET disk_used_bytes = $1 WHERE openid = $2",
                total, svc_openid
            )
        finally:
            await conn.close()
    except:
        pass

def _download_file_sync(encrypt_query_param: str, aes_key_b64: str, filename: str, user_id: str = "", file_type: str = "") -> str:
    """从 iLink CDN 下载并解密文件，返回保存路径（同步，用 urllib）"""
    if not encrypt_query_param or not aes_key_b64:
        return ""
    # 按用户分目录 + 按类型分子目录
    user_dir = FILE_DIR
    if user_id:
        safe_uid = re.sub(r'[^\w\.\-@]', '_', user_id.split('/')[-1])[:64]
        sub = _FILE_TYPE_DIRS.get(file_type, "other")
        user_dir = os.path.join(FILE_DIR, safe_uid, sub)
        os.makedirs(user_dir, exist_ok=True)
    url = f"https://novac2c.cdn.weixin.qq.com/c2c/download?encrypted_query_param={encrypt_query_param}"
    try:
        import urllib.request
        resp = urllib.request.urlopen(url, timeout=120)
        if resp.status != 200:
            return ""
        ciphertext = resp.read()
        # 解密 AES-128-ECB
        aes_key_hex = base64.b64decode(aes_key_b64).decode("ascii")
        aes_key = bytes.fromhex(aes_key_hex)
        from cryptography.hazmat.primitives.ciphers import Cipher, algorithms, modes
        from cryptography.hazmat.primitives import padding
        cipher = Cipher(algorithms.AES(aes_key), modes.ECB())
        decryptor = cipher.decryptor()
        padded = decryptor.update(ciphertext) + decryptor.finalize()
        unpadder = padding.PKCS7(128).unpadder()
        plaintext = unpadder.update(padded) + unpadder.finalize()
        # 写文件（带时间戳防重名）
        ts = int(time.time())
        safe_name = re.sub(r'[^\w\.\-]', '_', filename.split('/')[-1])
        if not safe_name or safe_name == '_':
            safe_name = f"file_{ts}"
        name, ext = os.path.splitext(safe_name)
        save_path = os.path.join(user_dir, f"{name}_{ts}{ext}")
        with open(save_path, "wb") as f:
            f.write(plaintext)
        return save_path
    except Exception:
        return ""

def _extract_pdf_text(path: str) -> str:
    """提取 PDF 文字"""
    try:
        import fitz
        doc = fitz.open(path)
        text = ""
        for page in doc:
            text += page.get_text() + "\n"
        doc.close()
        return text.strip()
    except:
        return "（PDF 文字提取失败）"

def _extract_pptx_text(path: str) -> str:
    """提取 PPTX 文字"""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text.strip() + "\n"
        return text.strip()
    except:
        return "（PPTX 文字提取失败）"

def _extract_text(msg: dict, from_user: str = "") -> str:
    """从消息中提取可处理的内容（文本/文件内容/语音识别结果）"""
    for item in msg.get("item_list", []):
        t = item.get("type")
        # 文本
        if t == ITEM_TEXT:
            return (item.get("text_item") or {}).get("text", "")
        
        # 文件（先下载归档，再识别内容）
        if t == ITEM_FILE:
            fi = item.get("file_item") or {}
            fname = fi.get("file_name", "") or fi.get("name", "") or f"file_{int(time.time())}"
            media = fi.get("media") or {}
            eq = media.get("encrypt_query_param", "")
            ak = media.get("aes_key", "")
            save_path = ""
            text = ""
            if eq and ak:
                save_path = _download_file_sync(eq, ak, fname, from_user, "file")
            if save_path:
                # 已归档到 downloads/ 目录
                # 同步到活跃项目
                try:
                    _PROJECT_STATE = None
                    _ps_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "..", "..", "mnt", "shared", "projects", ".project_state.json")
                    if not os.path.exists(_ps_path):
                        _ps_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "mnt", "shared", "projects", ".project_state.json")
                    if os.path.exists(_ps_path):
                        with open(_ps_path) as _ps_f:
                            _ps_data = json.load(_ps_f)
                        _active = _ps_data.get("active_project", "")
                        if _active:
                            _proj_dir = os.path.join(os.path.dirname(_ps_path), _active, "01_项目资料")
                            os.makedirs(_proj_dir, exist_ok=True)
                            _dst = os.path.join(_proj_dir, os.path.basename(save_path))
                            if not os.path.exists(_dst):
                                import shutil
                                shutil.copy2(save_path, _dst)
                except Exception:
                    pass
                try:
                    with open(save_path, "rb") as _f:
                        _header = _f.read(8)
                    if _header[:5] == b"%PDF-":
                        text = _extract_pdf_text(save_path)
                    elif _header[:4] == b"PK\x03\x04":
                        text = _extract_pptx_text(save_path)
                except:
                    pass
            size_str = ""
            if save_path:
                try:
                    sz = os.path.getsize(save_path)
                    size_str = f" ({sz // 1024}KB)"
                except:
                    pass
            if text:
                return f"[文件] {fname}{size_str}\n\n{text[:2000]}"
            if save_path:
                return f"[文件] {fname}{size_str} 已归档\n路径: {save_path}"
            return f"[文件] {fname}"
        
        # 图片（先下载归档）
        if t == ITEM_IMAGE:
            ii = item.get("image_item") or {}
            media = ii.get("media") or {}
            eq = media.get("encrypt_query_param", "")
            ak = media.get("aes_key", "")
            save_path = ""
            if eq and ak:
                save_path = _download_file_sync(eq, ak, f"image_{int(time.time())}.jpg", from_user, "image")
            if save_path:
                return f"[图片] {save_path}"
            return "[图片] (下载失败)"
        
        # 视频
        if t == ITEM_VIDEO:
            vi = item.get("video_item") or {}
            media = vi.get("media") or {}
            eq = media.get("encrypt_query_param", "")
            ak = media.get("aes_key", "")
            fname = vi.get("file_name", "") or f"video_{int(time.time())}"
            save_path = ""
            if eq and ak:
                save_path = _download_file_sync(eq, ak, fname, from_user, "video")
            if save_path:
                return f"[视频] {fname} 已归档"
            return f"[视频] {fname}"
        
        # 语音（先下载归档，再提取 ASR 文字）
        if t == ITEM_VOICE:
            voice = item.get("voice_item") or {}
            media = voice.get("media") or {}
            eq = media.get("encrypt_query_param", "")
            ak = media.get("aes_key", "")
            if eq and ak:
                _download_file_sync(eq, ak, f"voice_{int(time.time())}.silk", from_user, "voice")
            # 尝试 ASR 识别文字
            for key in ("asr_refer_text", "recog_text", "recognition", "text"):
                t = voice.get(key, "") or ""
                if t.strip():
                    return "[语音] " + t.strip()
            return "[语音] (识别失败，文件已归档)"
    
    return ""

# ══════════════════════════════════════════════
# HTTP API — 供 Agent Connector 调用来发消息
# ══════════════════════════════════════════════

from aiohttp import web

_running_bots: dict[str, BotSession] = {}  # bot_id → session

async def handle_send(request):
    """Agent Connector 调此接口发消息"""
    try:
        data = await request.json()
        bot_id = data.get("bot_id", "")
        to_user = data.get("to_user", "")
        text = data.get("text", "")
        ctx = data.get("context_token", "")

        bot = _running_bots.get(bot_id)
        if not bot:
            return web.json_response({"success": False, "error": "bot not found"}, status=404)

        # 检测 MEDIA: 路径 → 文件推送（支持嵌入文中的多条 MEDIA）
        media_paths = re.findall(r'MEDIA:(\S+?)(?:[\*\s）\)]|$)', text)
        if media_paths:
            clean = text
            ok_count = 0
            for mp in media_paths:
                mp_clean = re.sub(r'[\*\）\)]+$', '', mp)
                ok = await _send_file(bot.token, to_user, mp_clean, "", ctx, bot_id=bot_id)
                if ok:
                    ok_count += 1
                clean = clean.replace(f"MEDIA:{mp}", "").strip()
                await asyncio.sleep(0.5)
            clean = re.sub(r'\*+MEDIA:', '', clean)
            clean = re.sub(r'\*+（[^）]*[KkMmGg][Bb]）', '', clean)
            clean = re.sub(r'\*+', '', clean)
            clean = re.sub(r'\n{3,}', '\n\n', clean)
            clean = re.sub(r'[（(]\s*[)）]', '', clean)
            if clean.strip() and ok_count > 0:
                await send_text(bot.token, to_user, clean.strip(), ctx, bot_id=bot_id)
            return web.json_response({"success": ok_count > 0, "method": "media", "count": ok_count})

        await send_text(bot.token, to_user, text, ctx, bot_id=bot_id)
        return web.json_response({"success": True})
    except Exception as e:
        return web.json_response({"success": False, "error": str(e)}, status=500)

async def handle_reload(request):
    """重新加载 Bot 列表（扫码激活后触发的 reload，需同时启动 run_bot 任务）"""
    global _shutdown_event
    try:
        new_bots_raw = await load_bots()
        current_ids = {bid for bid, _ in new_bots_raw}
        running_ids = set(_running_tasks.keys())

        # 新 Bot：创建 run_bot 任务
        new_count = 0
        for bot_id, token in new_bots_raw:
            if bot_id not in running_ids:
                _bot_tokens[bot_id] = token
                t = asyncio.create_task(run_bot(bot_id, token, _shutdown_event))
                _running_tasks[bot_id] = t
                root_log.info("🆕 新 Bot（reload）: %s", bot_id)
                new_count += 1
            # 同步 _running_bots（供 HTTP API 用）
            if bot_id not in _running_bots:
                bot = BotSession(bot_id=bot_id, token=token, log=root_log)
                _running_bots[bot_id] = bot

        # 已停用的 Bot
        for bot_id in running_ids - current_ids:
            if bot_id in _running_tasks:
                _running_tasks[bot_id].cancel()
                del _running_tasks[bot_id]
            _bot_tokens.pop(bot_id, None)
            _running_bots.pop(bot_id, None)
            root_log.info("❌ Bot 已停用（reload）: %s", bot_id)

        return web.json_response({"success": True, "new": new_count, "total": len(new_bots_raw)})
    except Exception as e:
        root_log.error("reload error: %s", e)
        return web.json_response({"ok": False, "error": str(e)}, status=500)

async def handle_push_daily(request):
    """推送每日精读文章到所有订阅用户（daily_push=true）"""
    try:
        data = await request.json()
        title = data.get("title", "")
        link = data.get("link", "")
        sections = data.get("sections", [])

        # 构建卡片消息（纯文本，微信 iLink 不支持 markdown）
        lines = [f"🦞 享客虾评·每日精读 ｜ {title}", ""]
        for s in sections:
            emoji = s.get("emoji", "•")
            label = s.get("label", "")
            summary = s.get("summary", "")
            comment = s.get("comment", "")
            lines.append(f"{emoji} {label}")
            lines.append(f"  {summary}")
            if comment:
                lines.append(f"  🦞 {comment}")
            lines.append("")
        lines.append(f"🔗 {link}")
        lines.append("")
        lines.append("🦞 享客虾 · 每日精读")
        text = "\n".join(lines)

        import asyncpg
        conn = await asyncpg.connect(DB_DSN)
        try:
            rows = await conn.fetch("""
                SELECT s.openid, cb.channel_user_id
                FROM subscribers s
                JOIN channel_bindings cb ON cb.openid = s.openid
                WHERE s.daily_push = true
                  AND cb.is_active = true
                  AND cb.channel_type IN ('ilink', 'weixin')
            """)

            if not rows:
                root_log.info("[push-daily] 没有需要推送的订阅用户")
                return web.json_response({"success": True, "sent": 0, "total": 0})

            # 构建 user_id → bot 映射（从 _running_bots 按 bot_accounts.user_id 匹配）
            user_bot_map = {}
            import asyncpg as _apg2
            conn2 = await _apg2.connect(DB_DSN)
            try:
                bot_rows = await conn2.fetch(
                    "SELECT bot_id, user_id FROM bot_accounts WHERE is_active = true AND user_id IS NOT NULL AND user_id != ''"
                )
                for br in bot_rows:
                    bid = br["bot_id"]
                    uid = br["user_id"]
                    if bid in _running_bots and not is_paused(bid):
                        user_bot_map[uid] = _running_bots[bid]
            finally:
                await conn2.close()

            sent = 0
            failed = 0
            for row in rows:
                to_user = row["channel_user_id"]
                bot = user_bot_map.get(to_user)
                if not bot:
                    root_log.warning("[push-daily] 找不到 %s 对应的 bot，跳过", to_user[:20])
                    failed += 1
                    continue
                ok = await send_text(bot.token, to_user, text, "", bot_id=bot.bot_id)
                if ok:
                    sent += 1
                else:
                    failed += 1
                await asyncio.sleep(0.3)

            root_log.info(f"[push-daily] 推送完成: sent={sent}, failed={failed}, total={len(rows)}")
            return web.json_response({"success": True, "sent": sent, "failed": failed, "total": len(rows)})
        finally:
            await conn.close()
    except Exception as e:
        root_log.error(f"[push-daily] 异常: {e}")
        return web.json_response({"success": False, "error": str(e)}, status=500)


async def handle_health(request):
    return web.json_response({"ok": True, "bots": len(_running_bots), "uptime": time.time()})

async def handle_welcome(request):
    """外部调此接口推送欢迎消息"""
    try:
        data = await request.json()
        bot_id = data.get("bot_id", "")
        to_user = data.get("to_user", "")
        nickname = data.get("nickname", "")
        msg = data.get("message", f"✨ 欢迎你，{nickname}！\n\n享客虾 Bot 已就绪 ✅\n\n试试对我说「你好」开始对话，或发送「帮助」了解我能做什么。")
        ctx = data.get("context_token", "")

        bot = _running_bots.get(bot_id)
        if not bot:
            return web.json_response({"success": False, "error": "bot not found"}, status=404)

        await send_text(bot.token, to_user, msg, ctx, bot_id=bot_id)
        bot.welcome_sent.add(to_user)
        return web.json_response({"success": True})
    except Exception as e:
        return web.json_response({"success": False, "error": str(e)}, status=500)

async def handle_subscription_confirmed(request):
    """支付成功后，外部调此接口推送确认消息"""
    try:
        data = await request.json()
        bot_id = data.get("bot_id", "")
        to_user = data.get("to_user", "")
        nickname = data.get("nickname", "")
        if not nickname:
            nickname = "虾友"
        plan_name = data.get("plan_name", "会员")
        remain_days = data.get("remain_days", 30)
        xiake_points = data.get("xiake_points", 0)

        bot = _running_bots.get(bot_id)
        if not bot:
            # 没 bot_id 时，按 user_id 找 Bot
            for bid, b in _running_bots.items():
                if b.user_id and b.user_id.split("@")[0] == to_user.split("@")[0]:
                    bot = b
                    bot_id = bid
                    break
            if not bot:
                return web.json_response({"success": False, "error": "bot not found"}, status=404)

        msg = (
            f"🎉 开通成功！\n\n"
            f"🦞 {nickname}，欢迎成为享客虾 {plan_name} 伙伴！\n\n"
            f"📅 会员有效期 {remain_days} 天\n"
            f"🦞 虾点：{xiake_points} 点\n"
            f"✅ 解锁全部 AI 能力\n\n"
            f"回复「会员」查看会员状态\n"
            f"回复「续费」续期会员"
        )
        ctx = data.get("context_token", "")
        await send_text(bot.token, to_user, msg, ctx)

        # 同时清除该用户的首次欢迎标记（防止后续再发一遍欢迎）
        if bot_id in _running_bots:
            _running_bots[bot_id].welcome_sent.discard(to_user)

        root_log.info(f"[订阅确认] Bot={bot_id} → {to_user[:20]} ({nickname})")
        return web.json_response({"success": True})
    except Exception as e:
        root_log.error(f"[订阅确认] 异常: {e}")
        return web.json_response({"success": False, "error": str(e)}, status=500)
    """手动触发 Bot 重载（扫码激活后通知 keepalive 立即加载）"""
    try:
        bots = await load_bots()
        current_ids = {bid for bid, _ in bots}
        running_ids = set(_running_tasks.keys())

        new_count = 0
        # 新 Bot
        for bot_id, token in bots:
            if bot_id not in running_ids:
                _bot_tokens[bot_id] = token
                t = asyncio.create_task(run_bot(bot_id, token, asyncio.Event()))
                _running_tasks[bot_id] = t
                root_log.info("🆕 新 Bot（reload）: %s", bot_id)
                new_count += 1

        # 已停用
        for bot_id in running_ids - current_ids:
            _running_tasks[bot_id].cancel()
            del _running_tasks[bot_id]
            _bot_tokens.pop(bot_id, None)
            root_log.info("❌ Bot 已停用（reload）: %s", bot_id)

        return web.json_response({"success": True, "new": new_count, "total": len(bots)})
    except Exception as e:
        root_log.error("reload 异常: %s", e)
        return web.json_response({"success": False, "error": str(e)}, status=500)

async def handle_subscription_confirmed(request):
    """支付成功后，外部调此接口推送确认消息到用户 Bot"""
    try:
        data = await request.json()
        bot_id = data.get("bot_id", "")
        to_user = data.get("to_user", "")
        nickname = data.get("nickname", "虾友")
        plan_name = data.get("plan_name", "会员")
        remain_days = data.get("remain_days", 30)

        bot = _running_bots.get(bot_id)
        if not bot:
            for bid, b in _running_bots.items():
                if b.user_id and b.user_id.split("@")[0] == to_user.split("@")[0]:
                    bot = b
                    bot_id = bid
                    break
            if not bot:
                return web.json_response({"success": False, "error": "bot not found"}, status=404)

        msg = (
            f"🎉 开通成功！欢迎成为享客虾 {plan_name} 伙伴！\n\n"
            f"🦞 虾点：{data.get('xiake_points', 0)} 点\n"
            f"📅 到期：{data.get('expires_at', '')}（剩余 {remain_days} 天）\n\n"
            f"👉 续期会员：https://ai.pangoozn.com/subscribe"
        )
        ctx = data.get("context_token", "")
        await send_text(bot.token, to_user, msg, ctx, bot_id=bot_id)
        if bot_id in _running_bots:
            _running_bots[bot_id].welcome_sent.discard(to_user)
        root_log.info(f"[订阅确认] Bot={bot_id} → {to_user[:20]} ({nickname})")
        return web.json_response({"success": True})
    except Exception as e:
        root_log.error(f"[订阅确认] 异常: {e}")
        return web.json_response({"success": False, "error": str(e)}, status=500)

async def run_http_server():
    app = web.Application()
    app.router.add_post("/api/send", handle_send)
    app.router.add_post("/api/welcome", handle_welcome)
    app.router.add_post("/api/subscription-confirmed", handle_subscription_confirmed)
    app.router.add_post("/api/push-daily", handle_push_daily)
    app.router.add_get("/api/health", handle_health)
    app.router.add_post("/api/reload", handle_reload)

    runner = web.AppRunner(app)
    await runner.setup()
    site = web.TCPSite(runner, "127.0.0.1", HTTP_PORT)
    await site.start()
    root_log.info("🌐 HTTP API :9100 — send/welcome/health/reload")

# ══════════════════════════════════════════════
# DB 同步
# ══════════════════════════════════════════════

_running_tasks: dict[str, asyncio.Task] = {}
_bot_tokens: dict[str, str] = {}
_shutdown_event: asyncio.Event = None

async def sync_bots(shutdown: asyncio.Event):
    while not shutdown.is_set():
        try:
            bots = await load_bots()
            current_ids = {bid for bid, _ in bots}
            running_ids = set(_running_tasks.keys())

            # 新 Bot
            for bot_id, token in bots:
                if bot_id not in running_ids:
                    root_log.info("🆕 新 Bot: %s", bot_id)
                    _bot_tokens[bot_id] = token
                    t = asyncio.create_task(run_bot(bot_id, token, shutdown))
                    _running_tasks[bot_id] = t

            # Token 变更
            for bot_id, token in bots:
                old = _bot_tokens.get(bot_id)
                if old and old != token and bot_id in _running_tasks:
                    root_log.info("🔄 Token 变更，重启: %s", bot_id)
                    clear_pause(bot_id)
                    _running_tasks[bot_id].cancel()
                    _bot_tokens[bot_id] = token
                    await asyncio.sleep(2)
                    t = asyncio.create_task(run_bot(bot_id, token, shutdown))
                    _running_tasks[bot_id] = t

            # 已停用
            for bot_id in running_ids - current_ids:
                root_log.info("❌ Bot 已停用: %s", bot_id)
                _running_tasks[bot_id].cancel()
                del _running_tasks[bot_id]
                _bot_tokens.pop(bot_id, None)

            # 清理已完成
            for bot_id in list(_running_tasks.keys()):
                if _running_tasks[bot_id].done():
                    try:
                        _running_tasks[bot_id].result()
                    except asyncio.CancelledError:
                        pass
                    except Exception as e:
                        root_log.warning("Bot %s 异常退出: %s", bot_id, e)
                    del _running_tasks[bot_id]
                    _bot_tokens.pop(bot_id, None)

            # 同步 _running_bots (给 HTTP API 用)
            for bot_id, task in _running_tasks.items():
                if not task.done() and bot_id not in _running_bots:
                    # 还没跑起来，等等
                    pass
        except Exception as e:
            root_log.error("sync_bots 异常: %s", e)

        try:
            await asyncio.wait_for(asyncio.get_event_loop().create_future(), timeout=DB_RELOAD_INTERVAL)
        except asyncio.TimeoutError:
            pass
        except asyncio.CancelledError:
            break

# Bot session 注册 (run_bot 里调用)
def register_bot(bot_id: str, bot: BotSession):
    _running_bots[bot_id] = bot

def unregister_bot(bot_id: str):
    _running_bots.pop(bot_id, None)

# ══════════════════════════════════════════════
# Main
# ══════════════════════════════════════════════

async def main():
    root_log.info("=" * 50)
    root_log.info("Keepalive Service 启动")
    root_log.info("  保活层 :9100（暗号秒回 + iLink 长轮询）")
    root_log.info("  Agent 层 :9101（AI 推理服务）")
    root_log.info("  DB 刷新: %ds", DB_RELOAD_INTERVAL)
    root_log.info("=" * 50)

    _ensure_code_file()

    shutdown = asyncio.Event()
    global _shutdown_event
    _shutdown_event = shutdown

    # 信号处理
    loop = asyncio.get_event_loop()
    sig_count = 0
    def _signal():
        nonlocal sig_count
        sig_count += 1
        if sig_count >= 2:
            root_log.warning("二次信号，强制退出")
            sys.exit(1)
        root_log.info("📴 graceful shutdown...")
        shutdown.set()
    for s in (signal.SIGTERM, signal.SIGINT):
        try:
            loop.add_signal_handler(s, _signal)
        except NotImplementedError:
            pass

    # 启动 HTTP API
    await run_http_server()

    # 首次加载 Bot
    bots = await load_bots()
    root_log.info("首次加载 %d 个 Bot:", len(bots))
    for bot_id, token in bots:
        _bot_tokens[bot_id] = token
        t = asyncio.create_task(run_bot(bot_id, token, shutdown))
        _running_tasks[bot_id] = t
        root_log.info("  ✅ %s", bot_id)

    # 后台 DB 同步
    sync_task = asyncio.create_task(sync_bots(shutdown))

    root_log.info("✅ 全部就绪，等待消息...")
    await shutdown.wait()

    root_log.info("关闭中...")
    for task in list(_running_tasks.values()):
        task.cancel()
    sync_task.cancel()
    if _running_tasks or sync_task:
        await asyncio.gather(*_running_tasks.values(), sync_task, return_exceptions=True)
    root_log.info("Keepalive Service 已停止")

if __name__ == "__main__":
    try:
        asyncio.run(main())
    except KeyboardInterrupt:
        pass
